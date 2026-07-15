from datetime import datetime
import re


NAVY = "17365D"
GREEN = "4A896B"
LIME = "B8CC2E"
GOLD = "E5BC2D"
RED = "C5534F"
BLUE = "5388C4"
TEAL = "138A8A"
LIGHT_GREEN = "EAF2E1"
LIGHT_GRAY = "F2F4F5"
MID_GRAY = "C7CCD4"
DARK = "253238"
WHITE = "FFFFFF"
MUTED = "68757D"


CODE_LIKE_PATTERN = re.compile(
    r"(?:\{\s*['\"]?(?:gap_id|risk_id|affected_module|schema_version)['\"]?\s*:|"
    r"\[\s*\{\s*['\"]|<think>|```(?:json|python)?)",
    re.IGNORECASE,
)


def _plain_text(value):
    """Convert structured workflow values into presentation-safe prose."""
    if value is None:
        return ""
    if isinstance(value, str):
        text = value.strip()
        if (text.startswith("{") and text.endswith("}")) or (text.startswith("[") and text.endswith("]")):
            try:
                import ast

                return _plain_text(ast.literal_eval(text))
            except Exception:
                import re

                description = re.search(r"['\"]description['\"]\s*:\s*['\"]([^'\"]+)", text)
                if description:
                    return description.group(1)
        return text
    if isinstance(value, (int, float, bool)):
        return str(value)
    if isinstance(value, (list, tuple)):
        return "；".join(text for text in (_plain_text(item) for item in value) if text)
    if isinstance(value, dict):
        preferred = (
            "description", "summary", "thesis", "claim", "message", "text",
            "action", "reason", "expected_impact", "next_step", "risk",
        )
        for key in preferred:
            text = _plain_text(value.get(key))
            if text:
                return text
        ignored = {"schema_version", "gap_id", "risk_id", "insight_id", "cluster_id", "affected_module"}
        return "；".join(
            text for key, item in value.items()
            if key not in ignored and not key.endswith("_id")
            for text in [_plain_text(item)] if text
        )
    return str(value)


def _safe_filename(name):
    raw = name or "market_research_report.pptx"
    raw = re.sub(r"[^\w.\-]+", "_", raw, flags=re.UNICODE).strip("._")
    if not raw.lower().endswith(".pptx"):
        raw += ".pptx"
    return raw or "market_research_report.pptx"


def _shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text_frame.text or "").strip()


def _shape_font_sizes(shape, slide_height):
    text = _shape_text(shape)
    if not text:
        return []
    # Footer/source-note text is intentionally smaller and is not body copy.
    if shape.top >= int(slide_height * 0.91) or shape.height <= int(slide_height * 0.025):
        return []
    sizes = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(round(run.font.size.pt, 1))
    return sizes


def _rect_intersection_ratio(first, second):
    left = max(first.left, second.left)
    top = max(first.top, second.top)
    right = min(first.left + first.width, second.left + second.width)
    bottom = min(first.top + first.height, second.top + second.height)
    if right <= left or bottom <= top:
        return 0.0
    intersection = (right - left) * (bottom - top)
    smaller = min(first.width * first.height, second.width * second.height)
    return float(intersection / smaller) if smaller else 0.0


def _coverage_ratio(shapes, slide_width, slide_height, columns=40, rows=23):
    occupied = set()
    for shape in shapes:
        left = max(0, min(columns - 1, int(shape.left / slide_width * columns)))
        right = max(left, min(columns - 1, int((shape.left + shape.width) / slide_width * columns)))
        top = max(0, min(rows - 1, int(shape.top / slide_height * rows)))
        bottom = max(top, min(rows - 1, int((shape.top + shape.height) / slide_height * rows)))
        for x in range(left, right + 1):
            for y in range(top, bottom + 1):
                occupied.add((x, y))
    return round(len(occupied) / (columns * rows), 3)


def _audit_presentation(prs, slide_spec):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    spec_slides = slide_spec.get("slides") or []
    spec_by_page = {int(item.get("page_no") or index): item for index, item in enumerate(spec_slides, 1)}
    slide_checks = []
    overlap_pairs = 0
    out_of_bounds_shapes = 0
    code_like_text_hits = 0
    small_font_runs = 0
    content_font_sizes = []
    coverage_values = []

    for page_no, slide in enumerate(prs.slides, 1):
        meaningful = []
        page_overlap_pairs = 0
        page_out_of_bounds = 0
        page_code_hits = 0
        page_small_fonts = 0
        page_fonts = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            has_visual = bool(
                getattr(shape, "has_chart", False)
                or getattr(shape, "has_table", False)
                or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            if text or has_visual:
                meaningful.append(shape)
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                page_out_of_bounds += 1
            if text and CODE_LIKE_PATTERN.search(text):
                page_code_hits += 1
            sizes = _shape_font_sizes(shape, prs.slide_height)
            page_fonts.extend(sizes)
            page_small_fonts += sum(1 for size in sizes if size < 11)

        for index, first in enumerate(meaningful):
            for second in meaningful[index + 1:]:
                if _rect_intersection_ratio(first, second) >= 0.18:
                    page_overlap_pairs += 1

        coverage = _coverage_ratio(meaningful, prs.slide_width, prs.slide_height)
        coverage_values.append(coverage)
        overlap_pairs += page_overlap_pairs
        out_of_bounds_shapes += page_out_of_bounds
        code_like_text_hits += page_code_hits
        small_font_runs += page_small_fonts
        content_font_sizes.extend(page_fonts)

        spec = spec_by_page.get(page_no) or {}
        visual = spec.get("visual") or {}
        blocks = visual.get("blocks") or []
        chart_blocks = [item for item in blocks if str(item.get("type") or "").lower() == "chart"]
        table_blocks = [item for item in blocks if str(item.get("type") or "").lower() in {"table", "kpi_table", "heatmap"}]
        data_points = sum(
            len(series.get("values") or [])
            for block in chart_blocks
            for series in (block.get("series") or [])
            if isinstance(series, dict)
        )
        analysis_chars = sum(
            len(str(paragraph))
            for block in blocks
            if str(block.get("type") or "").lower() in {"analysis", "narrative", "action", "takeaway"}
            for paragraph in (block.get("paragraphs") or block.get("points") or [])
        )
        slide_checks.append({
            "slide_id": spec.get("slide_id") or f"S-{page_no:03d}",
            "page_no": page_no,
            "purpose": spec.get("purpose") or "unknown",
            "coverage_ratio": coverage,
            "overlap_pairs": page_overlap_pairs,
            "out_of_bounds_shapes": page_out_of_bounds,
            "code_like_text_hits": page_code_hits,
            "min_content_font_pt": min(page_fonts) if page_fonts else None,
            "small_font_runs": page_small_fonts,
            "chart_blocks": len(chart_blocks),
            "table_blocks": len(table_blocks),
            "data_points": data_points,
            "analysis_characters": analysis_chars,
        })

    finding_checks = [item for item in slide_checks if item["purpose"] == "finding"]
    finding_count = len(finding_checks)
    ratio = lambda count: round(count / finding_count, 3) if finding_count else 0.0
    layout_variants = {
        str((item.get("visual") or {}).get("layout_variant") or "")
        for item in spec_slides
        if (item.get("visual") or {}).get("layout_variant")
    }
    conclusion_titles = [
        item for item in spec_slides
        if item.get("purpose") in {"finding", "decision_summary", "recommendation"}
        and len(str(item.get("title") or "").strip()) >= 12
        and str(item.get("title") or "").strip() not in {"主要发现", "研究结论", "建议", "执行摘要"}
    ]
    decision_slides = [item for item in spec_slides if item.get("purpose") in {"finding", "decision_summary", "recommendation"}]
    metrics = {
        "mean_content_coverage_ratio": round(sum(coverage_values) / len(coverage_values), 3) if coverage_values else 0.0,
        "slides_in_coverage_range_ratio": round(sum(1 for value in coverage_values if 0.65 <= value <= 0.90) / len(coverage_values), 3) if coverage_values else 0.0,
        "min_content_font_pt": min(content_font_sizes) if content_font_sizes else None,
        "small_font_runs": small_font_runs,
        "overlap_pairs": overlap_pairs,
        "out_of_bounds_shapes": out_of_bounds_shapes,
        "code_like_text_hits": code_like_text_hits,
        "finding_slide_count": finding_count,
        "finding_slides_with_chart_ratio": ratio(sum(1 for item in finding_checks if 1 <= item["chart_blocks"] <= 2)),
        "finding_slides_with_tables": sum(1 for item in finding_checks if item["table_blocks"] > 0),
        "finding_slides_meeting_data_density_ratio": ratio(sum(1 for item in finding_checks if 8 <= item["data_points"] <= 24)),
        "finding_slides_meeting_analysis_depth_ratio": ratio(sum(1 for item in finding_checks if 160 <= item["analysis_characters"] <= 380)),
        "conclusion_title_ratio": round(len(conclusion_titles) / len(decision_slides), 3) if decision_slides else 0.0,
        "layout_variant_count": len(layout_variants),
    }
    warnings = []
    if overlap_pairs:
        warnings.append(f"Detected {overlap_pairs} potentially significant shape overlaps; review slide_checks.")
    if out_of_bounds_shapes:
        warnings.append(f"Detected {out_of_bounds_shapes} out-of-bounds shapes.")
    if small_font_runs:
        warnings.append(f"Detected {small_font_runs} content runs below 11pt.")
    if code_like_text_hits:
        warnings.append(f"Detected {code_like_text_hits} code-like visible text hits.")
    hard_fail = bool(out_of_bounds_shapes or code_like_text_hits)
    return {
        "schema_version": "ppt-visual-audit/1.0",
        "status": "fail" if hard_fail else ("pass_with_warnings" if warnings else "pass"),
        "metrics": metrics,
        "slide_checks": slide_checks,
        "warnings": warnings[:10],
    }


def _set_run(run, size, bold=False, color=DARK, font="Microsoft YaHei"):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _textbox(slide, x, y, w, h, text="", size=14, bold=False, color=DARK, fill=None, margin=0.08):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)
        shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = _plain_text(text)
    _set_run(run, size, bold, color)
    return shape


def _add_paragraphs(slide, items, x, y, w, h, size=12, color=DARK, bold_first=False):
    from pptx.util import Inches, Pt

    values = [_plain_text(value) for value in (items or [])]
    values = [value for value in values if value][:4]
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, value in enumerate(values):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = value
        paragraph.space_after = Pt(9)
        paragraph.line_spacing = 1.15
        for run in paragraph.runs:
            _set_run(run, size, bold_first and index == 0, color)
    return shape


def _add_header(slide, section, title, page_no):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor.from_string(GREEN)
    band.line.fill.background()
    _textbox(slide, 0.48, 0.22, 2.8, 0.24, section, 9, True, GREEN, margin=0)
    title_size = 23 if len(_plain_text(title)) <= 40 else 19
    _textbox(slide, 0.48, 0.48, 12.25, 0.88, title, title_size, True, NAVY, margin=0)
    _textbox(slide, 12.25, 7.12, 0.45, 0.18, page_no, 7, False, MUTED, margin=0)


def _add_source(slide, source_notes):
    if not source_notes:
        return
    notes = source_notes if isinstance(source_notes, list) else [source_notes]
    _textbox(slide, 0.5, 6.98, 11.45, 0.24, "来源：" + "；".join(_plain_text(note) for note in notes[:4]), 7, False, MUTED, margin=0)


def _body_parts(item):
    body = item.get("body") or {}
    if isinstance(body, list):
        return "", body, "", ""
    return (
        body.get("key_message") or item.get("headline") or "",
        body.get("supporting_points") or [],
        body.get("implication") or "",
        body.get("recommended_action") or "",
    )


def _add_bullets(slide, items, x=0.75, y=1.75, w=11.8, h=4.6, size=15):
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate([value for value in items if value][:8]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = _plain_text(item)
        paragraph.level = 0
        paragraph.space_after = 8
        for run in paragraph.runs:
            _set_run(run, size, False, DARK)
    return shape


def _chart_colors(chart):
    from pptx.dml.color import RGBColor

    colors = [BLUE, RED, GREEN, GOLD, TEAL, MID_GRAY]
    for index, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(colors[index % len(colors)])


def _add_chart(slide, visual, x=0.7, y=1.65, w=11.9, h=4.8, compact=False):
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
    from pptx.util import Inches, Pt

    categories = visual.get("categories") or []
    series = visual.get("series") or []
    valid_series = []
    for item in series[:6]:
        if not isinstance(item, dict):
            continue
        values = item.get("values") or []
        if len(values) != len(categories):
            continue
        valid_series.append(item)
    if not categories or not valid_series:
        return False

    chart_data = ChartData()
    chart_data.categories = [_plain_text(value) for value in categories]
    for item in valid_series:
        chart_data.add_series(_plain_text(item.get("name") or "Series"), [float(v or 0) for v in item.get("values")])

    chart_type = str(visual.get("chart_type") or "clustered_bar").lower()
    mapping = {
        "clustered_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "grouped_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "ranked_bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    }
    chart = slide.shapes.add_chart(
        mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
        Inches(x), Inches(y), Inches(w), Inches(h), chart_data,
    ).chart
    chart.has_legend = len(valid_series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8.5 if compact else 9.5)
    chart.has_title = False
    if chart_type != "pie":
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(8 if compact else 9)
        chart.category_axis.tick_labels.font.size = Pt(9.5 if compact else 10.5)
        if chart_type in {"ranked_bar", "stacked_bar"}:
            chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        values = [float(v or 0) for item in valid_series for v in item.get("values")]
        if values and min(values) >= 0 and max(values) <= 100 and visual.get("unit") in {"%", "percent", "percentage"}:
            chart.value_axis.minimum_scale = 0
            chart.value_axis.maximum_scale = 100
    point_count = len(categories) * len(valid_series)
    show_data_labels = visual.get("show_data_labels")
    if show_data_labels is None:
        show_data_labels = (
            (chart_type == "ranked_bar" and len(categories) <= 8)
            or (point_count <= 10)
            or (point_count <= 12 and w >= 6.5)
        )
    for plot in chart.plots:
        try:
            plot.gap_width = 65 if chart_type in {"clustered_bar", "grouped_bar"} else 45
        except Exception:
            pass
        plot.has_data_labels = bool(show_data_labels)
        if plot.has_data_labels:
            plot.data_labels.font.size = Pt(7.5 if compact else 8.5)
    _chart_colors(chart)
    return True


def _add_table(slide, visual, x=0.65, y=1.55, w=12.0, h=4.95, compact=False):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    table_data = visual.get("table") or {}
    headers = table_data.get("headers") or []
    rows = table_data.get("rows") or []
    if not headers or not rows:
        return False
    visible_rows = rows[:10 if compact else 12]
    shape = slide.shapes.add_table(
        len(visible_rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = shape.table
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = _plain_text(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(GREEN)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_run(run, 8 if compact else 10, True, WHITE)
    for row_index, row in enumerate(visible_rows, start=1):
        for column in range(len(headers)):
            value = row[column] if column < len(row) else ""
            cell = table.cell(row_index, column)
            cell.text = _plain_text(value)
            if row_index % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(LIGHT_GRAY)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_run(run, 7 if compact else 9, False, DARK)
    return True


def _add_panel(slide, panel, x, y, w, h, framed=True):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    if framed:
        border = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        border.fill.solid()
        border.fill.fore_color.rgb = RGBColor.from_string(WHITE)
        border.line.color.rgb = RGBColor.from_string("D8DDDF")
    title = _plain_text(panel.get("title") or panel.get("annotation") or "")
    if title:
        _textbox(slide, x + 0.12, y + 0.07, w - 0.24, 0.42, title, 13.5, True, NAVY, margin=0.02)
    panel_type = str(panel.get("type") or "chart").lower()
    content_y = y + (0.53 if title else 0.08)
    content_h = h - (0.61 if title else 0.16)
    if panel_type in {"analysis", "narrative", "takeaway", "action"}:
        paragraphs = panel.get("paragraphs") or panel.get("points") or [panel.get("headline"), panel.get("annotation")]
        if w > 10:
            paragraph_size = 14
        elif w < 4:
            paragraph_size = 11.5
        elif w < 5:
            paragraph_size = 12.5
        elif h < 2.7:
            paragraph_size = 12.5
        else:
            paragraph_size = 13.5
        _add_paragraphs(slide, paragraphs, x + 0.15, content_y + 0.04, w - 0.3, content_h - 0.08, paragraph_size)
        rendered = True
    elif panel_type in {"kpi_table", "table", "heatmap"}:
        rendered = _add_table(slide, panel, x + 0.08, content_y, w - 0.16, content_h, compact=True)
    else:
        rendered = _add_chart(slide, panel, x + 0.08, content_y, w - 0.16, content_h, compact=True)
    if not rendered:
        paragraphs = panel.get("paragraphs") or panel.get("points") or [panel.get("annotation")]
        _add_paragraphs(slide, paragraphs, x + 0.15, content_y + 0.05, w - 0.3, content_h - 0.1, 11.5)
    return True


def _add_story_layout(slide, visual):
    blocks = [block for block in (visual.get("blocks") or visual.get("panels") or []) if isinstance(block, dict)][:4]
    if not blocks:
        return False

    variant = str(visual.get("layout_variant") or "chart_left_analysis_right").lower()
    geometries = {
        "chart_left_analysis_right": [(0.5, 1.48, 7.25, 5.25), (7.98, 1.48, 4.85, 5.25)],
        "chart_right_analysis_left": [(0.5, 1.48, 4.75, 5.25), (5.48, 1.48, 7.35, 5.25)],
        "analysis_top_chart_bottom": [(0.5, 1.48, 12.33, 2.25), (0.5, 3.9, 12.33, 2.83)],
        "chart_top_analysis_bottom": [(0.5, 1.48, 12.33, 2.35), (0.5, 3.98, 12.33, 2.75)],
        "two_charts_analysis_right": [(0.5, 1.48, 7.35, 2.48), (0.5, 4.14, 7.35, 2.59), (8.08, 1.48, 4.75, 5.25)],
        "text_bands": [(0.5, 1.48, 12.33, 2.35), (0.5, 4.02, 12.33, 2.71)],
        "chart_focus_commentary": [(0.5, 1.48, 8.55, 5.25), (9.28, 1.48, 3.55, 5.25)],
        "balanced_columns": [(0.5, 1.48, 6.0, 5.25), (6.75, 1.48, 6.08, 5.25)],
    }
    boxes = geometries.get(variant, geometries["chart_left_analysis_right"])

    if len(blocks) > len(boxes):
        blocks = blocks[:len(boxes)]
    if len(blocks) < len(boxes):
        boxes = boxes[:len(blocks)]
    for block, box in zip(blocks, boxes):
        _add_panel(slide, block, *box, framed=False)
    return True


def _add_multi_panel(slide, visual):
    panels = [panel for panel in (visual.get("panels") or []) if isinstance(panel, dict)][:4]
    if not panels:
        return False
    if len(panels) == 1:
        boxes = [(0.55, 1.58, 12.23, 4.72)]
    elif len(panels) == 2:
        boxes = [(0.55, 1.58, 6.0, 4.72), (6.73, 1.58, 6.05, 4.72)]
    elif len(panels) == 3:
        boxes = [(0.5, 1.58, 4.02, 4.72), (4.66, 1.58, 4.02, 4.72), (8.82, 1.58, 4.02, 4.72)]
    else:
        boxes = [
            (0.55, 1.48, 6.0, 2.34), (6.73, 1.48, 6.05, 2.34),
            (0.55, 3.98, 6.0, 2.34), (6.73, 3.98, 6.05, 2.34),
        ]
    for panel, box in zip(panels, boxes):
        _add_panel(slide, panel, *box)
    return True


def _add_narrative_footer(slide, implication, action):
    if implication:
        _textbox(slide, 0.58, 6.31, 6.0, 0.6, "业务含义｜" + implication, 10.5, True, NAVY, LIGHT_GREEN, margin=0.06)
    if action:
        _textbox(slide, 6.75, 6.31, 6.0, 0.6, "建议行动｜" + action, 10.5, True, NAVY, LIGHT_GRAY, margin=0.06)


def _add_recommendation_plan(slide, visual, body):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    items = visual.get("items") or body.get("recommendations") or []
    items = [item for item in items if isinstance(item, dict)][:2]
    if not items:
        return False
    for index, item in enumerate(items):
        y = 1.52 + index * 2.62
        accent = slide.shapes.add_shape(1, Inches(0.52), Inches(y), Inches(0.08), Inches(2.35))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor.from_string(GREEN if index == 0 else GOLD)
        accent.line.fill.background()
        _textbox(slide, 0.72, y, 3.7, 0.3, f"优先级 {item.get('priority', index + 1)}｜{item.get('theme', '行动')}", 10, True, GREEN, margin=0)
        _textbox(slide, 0.72, y + 0.38, 3.85, 1.65, item.get("action") or "待明确行动", 15.5, True, NAVY, margin=0)
        rationale = _plain_text(item.get("reason") or item.get("evidence"))
        impact = _plain_text(item.get("expected_impact"))
        execution = "；".join(
            text for text in (
                "执行风险是" + _plain_text(item.get("risk")) if item.get("risk") else "",
                "下一步采用" + _plain_text(item.get("next_step") or item.get("validation")) if item.get("next_step") or item.get("validation") else "",
                "成功标准为" + _plain_text(item.get("success_metric")) if item.get("success_metric") else "",
            ) if text
        )
        paragraphs = []
        if rationale or impact:
            paragraphs.append((rationale + ("。预期将" + impact if impact else "")).rstrip("。") + "。")
        if execution:
            paragraphs.append(execution.rstrip("。") + "。")
        _add_paragraphs(slide, paragraphs, 4.72, y + 0.06, 8.05, 2.12, 12.5)
    return True


def _render_slide(prs, item, index):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout = str(item.get("layout") or "content").lower()
    title = item.get("title") or f"第 {index} 页"
    section = item.get("section") or "研究发现"
    key_message, points, implication, action = _body_parts(item)

    if layout == "cover":
        block = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor.from_string(NAVY)
        block.line.fill.background()
        _textbox(slide, 0.85, 1.45, 11.5, 1.1, title, 28, True, WHITE)
        _textbox(slide, 0.85, 2.62, 10.5, 0.65, key_message or "概念测试调研报告", 17, False, "DCE9CF")
        _textbox(slide, 0.85, 6.45, 4.0, 0.3, datetime.now().strftime("%Y"), 11, False, WHITE)
        return slide

    if layout == "section":
        block = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor.from_string(GREEN)
        block.line.fill.background()
        _textbox(slide, 0.9, 2.55, 11.5, 0.8, title, 31, True, WHITE)
        _textbox(slide, 0.9, 3.55, 10.5, 0.45, key_message, 16, False, WHITE)
        return slide

    _add_header(slide, section, title, item.get("page_no", index))
    visual = item.get("visual") or {}
    visual_type = str(visual.get("type") or "").lower()
    body = item.get("body") if isinstance(item.get("body"), dict) else {}
    rendered = False
    if layout == "recommendation" or visual_type == "recommendation_plan":
        rendered = _add_recommendation_plan(slide, visual, body)
    elif visual_type == "story_layout":
        rendered = _add_story_layout(slide, visual)
    elif visual_type == "multi_panel" or visual.get("panels"):
        rendered = _add_multi_panel(slide, visual)
    elif visual_type in {"kpi_comparison", "grouped_bar", "ranked_bar", "chart"}:
        rendered = _add_chart(slide, visual)
    elif visual_type in {"kpi_table", "table", "heatmap"}:
        rendered = _add_table(slide, visual)

    if rendered:
        if visual_type not in {"recommendation_plan", "story_layout"} and layout != "recommendation":
            _add_narrative_footer(slide, implication, action)
    else:
        if key_message:
            _textbox(slide, 0.65, 1.45, 12.0, 0.65, key_message, 15, True, GREEN)
        _add_bullets(slide, points or [implication, action], y=2.15, h=3.0)
        _add_narrative_footer(slide, implication, action)

    _add_source(slide, item.get("source_notes") or [])
    return slide


def render_market_research_pptx(payload, output_dir, reference_ppt, base_url):
    from pptx import Presentation

    slide_spec = payload.get("slide_spec") or {}
    slides = slide_spec.get("slides") or []
    if not isinstance(slides, list) or not slides:
        raise ValueError("slide_spec.slides must be a non-empty array")

    output = payload.get("output") or {}
    file_name = _safe_filename(output.get("file_name"))
    if file_name == "market_research_report.pptx":
        file_name = f"market_research_report_{datetime.now():%Y%m%d_%H%M%S}.pptx"
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / file_name

    prs = Presentation()
    if reference_ppt.exists():
        reference = Presentation(str(reference_ppt))
        prs.slide_width = reference.slide_width
        prs.slide_height = reference.slide_height

    for index, item in enumerate(slides, start=1):
        _render_slide(prs, item, index)

    prs.save(output_path)
    visual_audit = _audit_presentation(prs, slide_spec)
    return {
        "status": "success",
        "schema_version": "ppt-render-result/1.0",
        "file_url": f"{base_url}/files/{file_name}",
        "file_name": file_name,
        "file_path": str(output_path),
        "slide_count": len(slides),
        "message": "PPTX generated successfully with python-pptx.",
        "warnings": visual_audit["warnings"],
        "visual_audit": visual_audit,
    }
