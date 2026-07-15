import json
import sys
import threading
from pathlib import Path
from urllib.request import urlopen

import yaml
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
SERVICE_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SERVICE_DIR))

from server import Handler, ReportServer
from smoke_test import request_json, request_multipart

DIFY_STRING_LIMIT = 400000


def assert_string_limits(stage, outputs):
    for name, value in outputs.items():
        if isinstance(value, str) and len(value) >= DIFY_STRING_LIMIT:
            raise AssertionError(
                f"{stage}.{name} contains {len(value)} characters; limit is {DIFY_STRING_LIMIT - 1}"
            )


def load_code_node(nodes, node_id):
    namespace = {}
    exec(nodes[node_id]["data"]["code"], namespace)
    return namespace["main"]


def build_slide_spec(insight_id, evidence_ids, comparison):
    c1 = comparison["concept_values"]["C1"]["value"] or 0
    c2 = comparison["concept_values"]["C2"]["value"] or 0
    slides = []
    for page_no in range(1, 11):
        slide = {
            "slide_id": f"S-{page_no:03d}",
            "page_no": page_no,
            "section": "研究结论" if page_no > 1 else "封面",
            "purpose": "cover" if page_no == 1 else ("recommendation" if page_no == 10 else "finding"),
            "title": "DemoHealthNutritionLine新品概念测试" if page_no == 1 else f"经验证的研究结论 {page_no - 1}",
            "layout": "cover" if page_no == 1 else "content",
            "insight_ids": [] if page_no == 1 else [insight_id],
            "evidence_ids": [],
            "body": {
                "key_message": "概念测试调研报告" if page_no == 1 else "基于可追溯证据的业务判断",
                "supporting_points": [] if page_no == 1 else ["指标、受众和 Base 口径保持一致"],
                "implication": "用于产品与沟通决策" if page_no > 1 else "",
                "recommended_action": "继续验证核心卖点" if page_no == 10 else "",
            },
            "visual": {"type": "none"},
            "source_notes": [],
            "speaker_notes": "",
        }
        if page_no == 2:
            slide["evidence_ids"] = evidence_ids
            slide["visual"] = {
                "type": "kpi_comparison",
                "chart_type": "clustered_bar",
                "categories": ["概念 C1", "概念 C2"],
                "series": [{"name": comparison.get("metric_label") or "指标", "values": [c1, c2]}],
            }
            slide["source_notes"] = [comparison.get("question_title") or comparison.get("question_id")]
        slides.append(slide)
    return {"schema_version": "market-research-slides/1.0", "slides": slides}


def main():
    workflow_path = (
        Path(sys.argv[1]).resolve()
        if len(sys.argv) > 1
        else ROOT / "workflow" / "generic_market_research_ppt_workflow.yml"
    )
    document = yaml.safe_load(workflow_path.read_text(encoding="utf-8"))
    nodes = {node["id"]: node for node in document["workflow"]["graph"]["nodes"]}
    workbook = ROOT / "sample_input" / "synthetic_concept_test_data.xlsx"
    focus = "Focus on Concept A and Concept B purchase intent, clarity, value perception, and optimization direction"

    normalize_quality = load_code_node(nodes, "normalize_quality_score")
    perfect_dimensions = {
        "evidence_accuracy": {"score": 15},
        "research_methodology": {"score": 15},
        "insight_depth": {"score": 25},
        "recommendation_actionability": {"score": 15},
        "narrative_coherence": {"score": 10},
        "chart_effectiveness": {"score": 10},
        "layout_readability": {"score": 7},
        "delivery_integrity": {"score": 3},
    }
    quality_result = normalize_quality(json.dumps({
        "schema_version": "quality-evaluation/1.0",
        "dimensions": perfect_dimensions,
        "hard_gate_violations": [],
    }, ensure_ascii=False))
    if quality_result["quality_score"] != 100.0 or quality_result["quality_grade"] != "excellent":
        raise AssertionError("quality score normalizer did not preserve the fixed 100-point rubric")
    gated_result = normalize_quality(json.dumps({
        "schema_version": "quality-evaluation/1.0",
        "dimensions": perfect_dimensions,
        "hard_gate_violations": ["visible code-like text"],
    }, ensure_ascii=False))
    if gated_result["quality_score"] != 59.0 or gated_result["hard_gate_status"] != "hard_gate_failed":
        raise AssertionError("quality score hard gate was not applied")

    server = ReportServer(("127.0.0.1", 0), Handler)
    port = server.server_address[1]
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    base_url = f"http://127.0.0.1:{port}"

    try:
        extracted = request_multipart(
            base_url + "/extract-market-data",
            fields={"business_focus": focus, "max_evidence": 600},
            file_field="data_file",
            file_path=workbook,
        )

        normalize = load_code_node(nodes, "normalize_research_data")
        normalized = normalize(json.dumps(extracted, ensure_ascii=False))
        assert_string_limits("normalize_research_data", normalized)
        build_matrix = load_code_node(nodes, "build_analysis_matrix")
        matrix_result = build_matrix(normalized["research_evidence"], focus)
        assert_string_limits("build_analysis_matrix", matrix_result)
        matrix = json.loads(matrix_result["analysis_matrix"])
        catalog = json.loads(matrix_result["evidence_catalog"])
        significant_comparisons = [item for item in matrix["comparisons"] if item.get("significant")]
        non_significant_comparisons = [item for item in matrix["comparisons"] if not item.get("significant")]
        selected_comparisons = (significant_comparisons + non_significant_comparisons)[:8]
        if len(selected_comparisons) < 8:
            raise AssertionError("expected at least eight comparable matrix items for insight validation")
        comparison = selected_comparisons[0]
        evidence_ids = comparison["evidence_ids"]
        insights = []
        for index, item in enumerate(selected_comparisons, start=1):
            insight_id = f"INS-TEST-{index:03d}"
            insights.append({
                "insight_id": insight_id,
                "priority": "high" if index <= 2 else "medium",
                "type": "concept_comparison",
                "audience": item.get("audience") or "All",
                "comparison_id": item["comparison_id"],
                "claim": (
                    f"{item.get('question_title') or item.get('question_id')}在两款概念间显著不同"
                    if item.get("significant")
                    else f"{item.get('question_title') or item.get('question_id')}呈方向性差异，但未达统计显著"
                ),
                "business_implication": "该指标可用于判断概念表达与目标受众需求的匹配程度",
                "recommended_action": "围绕优势表达继续验证购买意愿与价值感",
                "evidence_ids": item["evidence_ids"],
                "comparison": {
                    "direction": "C2_over_C1" if item.get("delta_c2_vs_c1", 0) > 0 else "C1_over_C2",
                    "statistically_significant": bool(item.get("significant")),
                },
                "counter_evidence_ids": [],
                "confidence": "medium",
            })

        insight_id = insights[0]["insight_id"]
        insight_package = {
            "schema_version": "research-insights/1.0",
            "decision": {
                "status": "conditional_go",
                "recommended_concept_id": "C1",
                "confidence": "medium",
                "rationale_insight_ids": [item["insight_id"] for item in insights[:2]],
            },
            "executive_summary": [
                "核心概念指标已按统一口径完成比较",
                "结论均保留可追溯证据编号",
                "建议围绕购买意愿和理解度继续优化",
            ],
            "insights": insights,
            "segment_opportunities": [],
            "risks": [{
                "risk_id": "RISK-001",
                "description": "部分结果需要结合样本质量谨慎解释",
                "severity": "medium",
                "mitigation": "在下一轮验证中复核关键指标",
                "evidence_ids": evidence_ids,
            }],
            "recommendations": [
                {
                    "priority": 1,
                    "theme": "沟通",
                    "action": "简化核心概念表达",
                    "reason": "提升理解度并强化购买理由",
                    "insight_ids": [insights[0]["insight_id"]],
                    "risk": "过度简化可能削弱差异化",
                    "next_step": "开展两版概念文案快速验证",
                },
                {
                    "priority": 2,
                    "theme": "下一步研究",
                    "action": "复核重点指标",
                    "reason": "确认结论在目标受众中的稳定性",
                    "insight_ids": [],
                    "risk": "样本不足会影响判断",
                    "next_step": "补充目标人群样本并复测",
                },
            ],
            "data_gaps": [
                {
                    "gap_id": "GAP-TEST-001",
                    "description": "缺少价格弹性数据，当前定价建议只依据价值感与购买意愿。",
                    "affected_module": "commercial_viability",
                }
            ],
        }
        validate_insights = load_code_node(nodes, "validate_insights")
        insight_result = validate_insights(
            "<think>internal reasoning must be ignored</think>" + json.dumps(insight_package, ensure_ascii=False),
            matrix_result["analysis_matrix"],
            matrix_result["evidence_catalog"],
        )
        assert_string_limits("validate_insights", insight_result)
        if insight_result["valid"] != "true":
            raise AssertionError(insight_result["validation_error"])
        truncated_result = validate_insights(
            "<think>unfinished reasoning</think>{\"schema_version\":",
            matrix_result["analysis_matrix"],
            matrix_result["evidence_catalog"],
        )
        if truncated_result["valid"] != "true" or "fallback" not in truncated_result["validation_error"]:
            raise AssertionError("truncated insight output did not receive a usable fallback")

        slide_spec = build_slide_spec(insight_id, evidence_ids, comparison)
        validate_slides = load_code_node(nodes, "validate_slide_spec")
        slide_result = validate_slides(
            "<think>ignored</think>" + json.dumps(slide_spec, ensure_ascii=False),
            insight_result["clean_insights"],
            matrix_result["evidence_catalog"],
            matrix_result["analysis_matrix"],
        )
        assert_string_limits("validate_slide_spec", slide_result)
        if slide_result["valid"] != "true":
            raise AssertionError(slide_result["validation_error"])
        fallback_slide_result = validate_slides(
            "<think>unfinished</think>{\"schema_version\":",
            insight_result["clean_insights"],
            matrix_result["evidence_catalog"],
            matrix_result["analysis_matrix"],
        )
        fallback_slide_spec = json.loads(fallback_slide_result["clean_slide_spec"])
        fallback_slides = fallback_slide_spec.get("slides", [])
        if fallback_slide_result["valid"] != "true" or not 12 <= len(fallback_slides) <= 16:
            raise AssertionError(
                f"invalid slide output produced {len(fallback_slides)} slides: "
                f"{fallback_slide_result['validation_error']}"
            )
        finding_slides = [slide for slide in fallback_slides if slide.get("purpose") == "finding"]
        if not finding_slides:
            raise AssertionError("fallback did not create finding slides")
        layout_variants = {
            (slide.get("visual") or {}).get("layout_variant")
            for slide in fallback_slides
            if (slide.get("visual") or {}).get("type") == "story_layout"
        }
        if len(layout_variants) < 3:
            raise AssertionError("fallback did not create enough visual layout variety")
        for slide in finding_slides:
            visual = slide.get("visual") or {}
            blocks = visual.get("blocks") or []
            block_types = [str(block.get("type") or "").lower() for block in blocks]
            if visual.get("type") != "story_layout" or not 2 <= len(blocks) <= 3:
                raise AssertionError("fallback finding slide did not use variable story layouts")
            if not 1 <= block_types.count("chart") <= 2:
                raise AssertionError("fallback finding slide did not create a focused chart story")
            if sum(kind in {"analysis", "narrative", "takeaway", "action"} for kind in block_types) < 1:
                raise AssertionError("fallback finding slide did not create an analysis block")
            if any(kind in {"table", "kpi_table", "heatmap"} for kind in block_types):
                raise AssertionError("fallback finding slide still contains a table")
            narrative_chars = sum(
                len(str(paragraph))
                for block in blocks
                if str(block.get("type") or "").lower() in {"analysis", "narrative", "takeaway", "action"}
                for paragraph in (block.get("paragraphs") or block.get("points") or [])
            )
            if narrative_chars < 180:
                raise AssertionError(
                    f"{slide.get('slide_id')} fallback finding slide contains only "
                    f"{narrative_chars} paragraph characters"
                )
        serialized_slides = json.dumps(fallback_slide_spec, ensure_ascii=False)
        if "gap_id" in serialized_slides or "affected_module" in serialized_slides:
            raise AssertionError("structured gap metadata leaked into visible slide content")
        recommendation_slides = [slide for slide in fallback_slides if slide.get("purpose") == "recommendation"]
        if len(recommendation_slides) < 2:
            raise AssertionError("fallback did not create two recommendation slides")
        for slide in recommendation_slides:
            for item in (slide.get("visual") or {}).get("items", []):
                if len(str(item.get("reason") or "")) < 80:
                    raise AssertionError("recommendation rationale is too short")

        qa_package = {
            "schema_version": "content-qa/1.0",
            "overall_status": "pass",
            "checks": [],
            "blockers": [],
            "warnings": [],
            "manual_review_items": [],
        }
        validate_qa = load_code_node(nodes, "validate_content_qa")
        qa_result = validate_qa("<think>unfinished</think>{\"schema_version\":")
        assert_string_limits("validate_content_qa", qa_result)
        if qa_result["valid"] != "true":
            raise AssertionError(qa_result["validation_error"])

        rendered = request_json(
            base_url + "/render-pptx",
            {
                "slide_spec": fallback_slide_spec,
                "evidence_registry": catalog,
                "output": {"file_name": "generated_demo_market_research_report.pptx"},
            },
        )
        visual_audit = rendered.get("visual_audit") or {}
        if visual_audit.get("schema_version") != "ppt-visual-audit/1.0":
            raise AssertionError("render response is missing ppt-visual-audit/1.0")
        audit_metrics = visual_audit.get("metrics") or {}
        if audit_metrics.get("code_like_text_hits") != 0:
            raise AssertionError("visual audit detected code-like text in rendered PPTX")
        if len(visual_audit.get("slide_checks") or []) != len(fallback_slides):
            raise AssertionError("visual audit did not return one check per slide")
        prepare_download = load_code_node(nodes, "prepare_ppt_download")
        prepared_download = prepare_download(json.dumps(rendered, ensure_ascii=False))
        with urlopen(prepared_download["download_url"], timeout=30) as response:
            downloaded_pptx = response.read()
            content_type = response.headers.get_content_type()
            content_disposition = response.headers.get("Content-Disposition", "")
        if not downloaded_pptx.startswith(b"PK"):
            raise AssertionError("download endpoint did not return a PPTX/ZIP binary")
        if content_type != "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            raise AssertionError(f"unexpected PPTX content type: {content_type}")
        if "attachment" not in content_disposition.lower():
            raise AssertionError("download endpoint did not return an attachment header")
        output_path = Path(rendered["file_path"])
        presentation = Presentation(str(output_path))
        if len(presentation.slides) != len(fallback_slides):
            raise AssertionError(f"expected {len(fallback_slides)} rendered slides, got {len(presentation.slides)}")
        rendered_text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        if "{'gap_id'" in rendered_text or "affected_module" in rendered_text:
            raise AssertionError("code-like dictionary text leaked into the rendered PPTX")

        result = {
            "status": "success",
            "workflow": str(workflow_path),
            "transport": "multipart/form-data",
            "extracted_evidence": len(json.loads(normalized["research_evidence"])["evidence_registry"]),
            "comparisons": len(matrix["comparisons"]),
            "insight_validation": insight_result["valid"],
            "slide_validation": slide_result["valid"],
            "content_qa_validation": qa_result["valid"],
            "quality_score_normalizer": quality_result["quality_score"],
            "quality_hard_gate_cap": gated_result["quality_score"],
            "insight_fallback_diagnostic": truncated_result["validation_error"],
            "slide_fallback_diagnostic": fallback_slide_result["validation_error"],
            "rendered_slides": len(presentation.slides),
            "visual_audit_status": visual_audit.get("status"),
            "visual_audit_metrics": audit_metrics,
            "downloaded_bytes": len(downloaded_pptx),
            "download_content_type": content_type,
            "download_disposition": content_disposition,
            "output_characters": {
                "research_evidence": len(normalized["research_evidence"]),
                "analysis_matrix": len(matrix_result["analysis_matrix"]),
                "evidence_catalog": len(matrix_result["evidence_catalog"]),
                "clean_insights": len(insight_result["clean_insights"]),
                "clean_slide_spec": len(slide_result["clean_slide_spec"]),
                "clean_qa": len(qa_result["clean_qa"]),
            },
            "output_file": str(output_path),
        }
        print(json.dumps(result, ensure_ascii=False, indent=2))
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=5)


if __name__ == "__main__":
    main()
