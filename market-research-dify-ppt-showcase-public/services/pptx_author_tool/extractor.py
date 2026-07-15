import base64
import gzip
import hashlib
import json
import re
import tempfile
from collections import defaultdict
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import Request, urlopen


QUESTION_RE = re.compile(r"^\s*\[([^\]]+)\](.*)$")
VALUE_RE = re.compile(r"^\s*(-?\d+(?:\.\d+)?)\s*([A-Z]+)?(\*{1,2})?\s*$")
DATASET_ID_RE = re.compile(r"^DS-[A-F0-9]{24}$")
MAX_EVIDENCE_LIMIT = 600
DIFY_SAFE_OUTPUT_CHARS = 360000
KEY_QUESTION_PREFIXES = (
    "Q1", "Q2", "Q3", "Q4", "Q5", "Q6", "Q7", "Q8", "Q9", "Q10",
    "Q11", "Q12", "Q13", "Q15", "Q17", "Q18", "Q19", "Q22", "Q23",
)
SUMMARY_ROWS = {"T1B", "T2B", "B1B", "B2B", "Mean", "Average", "Definitely would buy", "Probably would buy"}
CONCEPT_ALIASES = {
    "C1": ("Concept A",),
    "C2": ("Concept B",),
    "EXISTING_POWDER": ("Existing Powder",),
    "WOMEN_MULTI": ("Women Multi",),
    "MEN_MULTI": ("Men Multi",),
    "DOUBLE_X": ("Existing Product X",),
}
QUESTION_TITLES = {
    "Q1": "Overall appeal",
    "Q2": "Reasons for appeal",
    "Q3": "Reasons for rejection",
    "Q4-1": "Purchase intent - Partner",
    "Q4-2": "Purchase intent - Consumer",
    "Q5": "Relevance",
    "Q6": "Uniqueness",
    "Q7": "Ease of understanding",
    "Q8": "Credibility",
    "Q9": "Value perception",
    "Q10a": "Message appeal",
    "Q11": "Recommendation intent",
    "Q12": "Reasons for non-recommendation",
    "Q13": "Purchase purpose",
    "Q15c": "Product-solution fit",
    "Q17": "Impact on current routine",
    "Q18a": "Positioning",
    "Q19a": "Reasons to increase current product usage",
    "Q19b": "Reasons to maintain current product usage",
    "Q19c": "Reasons to reduce current product usage",
    "Q19d": "Reasons to replace current product usage",
    "Q22a": "Ease of selling",
    "Q22b": "Selling barriers",
    "Q23": "Recommended audience",
}


def _jsonish(value):
    if isinstance(value, dict):
        return value
    if isinstance(value, str):
        text = value.strip()
        if text.startswith("{"):
            try:
                parsed = json.loads(text)
                return parsed if isinstance(parsed, dict) else {"url": text}
            except json.JSONDecodeError:
                return {"url": text}
        return {"url": text}
    return {}


def materialize_input_file(file_ref, suffix=".xlsx"):
    ref = _jsonish(file_ref)
    raw_bytes = ref.get("content_bytes")
    encoded = ref.get("content_base64") or ref.get("base64")
    if isinstance(raw_bytes, bytes):
        data = raw_bytes
    elif encoded:
        data = base64.b64decode(encoded)
    else:
        location = ref.get("url") or ref.get("remote_url") or ref.get("path")
        if not location:
            raise ValueError("data_file must contain url, path, or content_base64")
        parsed = urlparse(str(location))
        if parsed.scheme in {"http", "https"}:
            request = Request(str(location), headers={"User-Agent": "Dify-Market-Research-Tool/1.0"})
            with urlopen(request, timeout=120) as response:
                data = response.read()
        else:
            path = Path(str(location)).expanduser().resolve()
            if not path.exists():
                raise ValueError(f"input file does not exist: {path}")
            data = path.read_bytes()

    if not data:
        raise ValueError("uploaded data_file is empty")

    handle = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
    handle.write(data)
    handle.close()
    return Path(handle.name)


def _normalize_text(value):
    if value is None:
        return ""
    return re.sub(r"\s+", " ", str(value)).strip()


def _fill_forward(values):
    output = []
    last = ""
    for value in values:
        text = _normalize_text(value)
        if text:
            last = text
        output.append(last)
    return output


def _parse_number(value):
    if value is None:
        return None, "", [], ""
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return float(value), str(value), [], ""
    display = _normalize_text(value)
    if display in {"", "-", "–", "—", ">", "<"}:
        return None, display, [], ""
    match = VALUE_RE.match(display.replace("%", ""))
    if not match:
        return None, display, [], ""
    number = float(match.group(1))
    letters = list(match.group(2) or "")
    stars = match.group(3) or ""
    return number, display, letters, stars


def _base_quality(base_value):
    _, display, _, stars = _parse_number(base_value)
    if "**" in display or stars == "**":
        return "very_small"
    if "*" in display or stars == "*":
        return "small"
    return "adequate"


def _concept_id(text, sheet_name):
    haystack = f"{text} {sheet_name}"
    for concept_id, aliases in CONCEPT_ALIASES.items():
        if any(alias in haystack for alias in aliases):
            return concept_id
    sheet_map = {
        "Header1_1": "C1", "Header1_2": "C2", "Header1_3": "C1",
        "Header1_4": "C2", "Header1_5": "C1", "Header1_6": "C2",
    }
    return sheet_map.get(sheet_name, "UNSPECIFIED")


def _audience(question_title, header_text):
    text = f"{question_title} {header_text}"
    if "消费者" in text:
        return "Consumer"
    if "Partner" in text or "channel partner" in text.lower():
        return "Partner"
    return "All"


def _segment(header_path):
    clean = []
    for item in header_path:
        if not item or item.startswith("("):
            continue
        if any(alias in item for aliases in CONCEPT_ALIASES.values() for alias in aliases):
            continue
        if item in {"概念", "组别出示-First Read", "组别出示-Second Read", "%", "总体"}:
            continue
        if item not in clean:
            clean.append(item)
    if not clean:
        return {"dimension": "total", "value": "total"}
    if len(clean) == 1:
        return {"dimension": "segment", "value": clean[0]}
    return {"dimension": clean[-2], "value": clean[-1]}


def _metric_id(question_id, row_label):
    digest = hashlib.sha1(f"{question_id}|{row_label}".encode("utf-8")).hexdigest()[:10]
    return f"MET-{question_id.replace('-', '_')}-{digest}"


def _evidence_id(sheet_name, question_id, row_index, column_index):
    raw = f"{sheet_name}|{question_id}|{row_index}|{column_index}"
    return "EV-" + hashlib.sha1(raw.encode("utf-8")).hexdigest()[:14].upper()


def _priority(question_id, row_label, segment, concept_id):
    score = 10
    if question_id.startswith(KEY_QUESTION_PREFIXES):
        score += 40
    if row_label in SUMMARY_ROWS:
        score += 35
    if segment.get("dimension") == "total":
        score += 25
    if concept_id in {"C1", "C2"}:
        score += 20
    if row_label in {"Total", "S.D.", "标准差"}:
        score -= 50
    return score


def _focus_terms(text):
    parts = re.split(r"[，,。；;、\n\s]+", text or "")
    return [part for part in parts if len(part) >= 2][:20]


def _dataset_id(path):
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for chunk in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(chunk)
    return "DS-" + digest.hexdigest()[:24].upper()


def _dataset_path(dataset_dir, dataset_id):
    if not DATASET_ID_RE.fullmatch(str(dataset_id or "")):
        raise ValueError("invalid dataset_id")
    return Path(dataset_dir) / f"{dataset_id}.json.gz"


def _write_dataset(dataset_dir, dataset_id, package):
    path = _dataset_path(dataset_dir, dataset_id)
    path.parent.mkdir(parents=True, exist_ok=True)
    with gzip.open(path, "wt", encoding="utf-8", compresslevel=6) as target:
        json.dump(package, target, ensure_ascii=False, separators=(",", ":"))
    return path


def _read_dataset(dataset_dir, dataset_id):
    path = _dataset_path(dataset_dir, dataset_id)
    if not path.exists():
        raise ValueError(f"dataset not found: {dataset_id}")
    with gzip.open(path, "rt", encoding="utf-8") as source:
        return json.load(source)


def _parse_sheet(ws):
    rows = list(ws.iter_rows(values_only=True))
    question_starts = []
    for index, row in enumerate(rows):
        first = row[0] if row else None
        if isinstance(first, str) and QUESTION_RE.match(first):
            question_starts.append(index)

    evidence = []
    metric_definitions = {}
    table_summaries = []

    for position, start in enumerate(question_starts):
        end = question_starts[position + 1] if position + 1 < len(question_starts) else len(rows)
        match = QUESTION_RE.match(str(rows[start][0]))
        question_id = match.group(1).strip()
        question_title = QUESTION_TITLES.get(question_id, match.group(2).strip())

        base_index = None
        for index in range(start + 1, min(end, start + 18)):
            label = _normalize_text(rows[index][0] if rows[index] else None)
            if label.startswith("Base"):
                base_index = index
                break
        if base_index is None or base_index + 1 >= end:
            continue

        max_columns = max(len(row) for row in rows[start:end])
        header_rows = []
        for header_index in range(start + 1, base_index):
            padded = list(rows[header_index]) + [None] * (max_columns - len(rows[header_index]))
            header_rows.append(_fill_forward(padded))

        unit_row = rows[base_index + 1]
        bases = rows[base_index]
        block_count = 0

        for row_index in range(base_index + 2, end):
            row = rows[row_index]
            row_label = _normalize_text(row[0] if row else None)
            if not row_label or row_label.startswith(("Proportions/Means", "Overlap formulae", "Significance")):
                continue
            if row_label in {"Total", "S.D.", "标准差"}:
                continue

            for column_index in range(1, min(len(row), max_columns)):
                value_raw, value_display, significance, _ = _parse_number(row[column_index])
                if value_raw is None:
                    continue

                header_path = []
                for header in header_rows:
                    text = _normalize_text(header[column_index])
                    if text and text not in header_path:
                        header_path.append(text)
                header_text = " | ".join(header_path)
                concept_id = _concept_id(header_text, ws.title)
                segment = _segment(header_path)
                base_value = bases[column_index] if column_index < len(bases) else None
                base_number, base_display, _, _ = _parse_number(base_value)
                unit = _normalize_text(unit_row[column_index] if column_index < len(unit_row) else "")
                normalized_unit = "number" if row_label == "Mean" else ("percent" if unit == "%" else (unit or "number"))
                metric_id = _metric_id(question_id, row_label)

                metric_definitions[metric_id] = {
                    "metric_id": metric_id,
                    "question_id": question_id,
                    "question_title": question_title,
                    "label": row_label,
                    "unit": normalized_unit,
                    "summary_method": row_label if row_label in SUMMARY_ROWS else "response_option",
                }

                item = {
                    "evidence_id": _evidence_id(ws.title, question_id, row_index + 1, column_index + 1),
                    "question_id": question_id,
                    "question_title": question_title,
                    "audience": _audience(question_title, header_text),
                    "concept_id": concept_id,
                    "segment": segment,
                    "metric_id": metric_id,
                    "metric_label": row_label,
                    "value_raw": value_raw,
                    "value_display": value_display + ("%" if normalized_unit == "percent" and "%" not in value_display else ""),
                    "unit": normalized_unit,
                    "base": int(base_number) if base_number is not None else None,
                    "base_display": base_display,
                    "base_quality": _base_quality(base_value),
                    "significance": {
                        "tested": bool(significance),
                        "level": 0.05 if significance else None,
                        "different_from_columns": significance,
                    },
                    "source": {
                        "sheet": ws.title,
                        "row": row_index + 1,
                        "column": column_index + 1,
                        "question": f"[{question_id}]{question_title}",
                        "header": header_text,
                    },
                }
                item["priority_score"] = _priority(question_id, row_label, segment, concept_id)
                evidence.append(item)
                block_count += 1

        table_summaries.append({
            "sheet": ws.title,
            "question_id": question_id,
            "question_title": question_title,
            "evidence_count": block_count,
        })

    return evidence, list(metric_definitions.values()), table_summaries


def _reference_outline(reference_ppt):
    if not reference_ppt.exists():
        return []
    from pptx import Presentation

    prs = Presentation(str(reference_ppt))
    outline = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _normalize_text(shape.text)
                if text and text not in texts:
                    texts.append(text)
        outline.append({
            "page_no": index,
            "title": (texts[0] if texts else "")[:240],
            "key_texts": [text[:240] for text in texts[1:4]],
        })
    return outline


def extract_market_research(payload, reference_ppt):
    from openpyxl import load_workbook

    local_file = materialize_input_file(payload.get("data_file"), suffix=".xlsx")
    workbook = None
    try:
        workbook = load_workbook(local_file, read_only=True, data_only=True)
        all_evidence = []
        metric_map = {}
        tables = []
        for worksheet in workbook.worksheets:
            if worksheet.title == "TOC":
                continue
            evidence, metrics, summaries = _parse_sheet(worksheet)
            all_evidence.extend(evidence)
            tables.extend(summaries)
            for metric in metrics:
                metric_map[metric["metric_id"]] = metric

        deduplicated = {}
        for item in all_evidence:
            signature = (
                item["question_id"], item["metric_label"], item["audience"],
                item["concept_id"], item["segment"]["dimension"],
                item["segment"]["value"], item["value_raw"], item["base"],
            )
            current = deduplicated.get(signature)
            if current is None or item["priority_score"] > current["priority_score"]:
                deduplicated[signature] = item

        requested_max_evidence = int(payload.get("max_evidence") or MAX_EVIDENCE_LIMIT)
        max_evidence = max(100, min(requested_max_evidence, MAX_EVIDENCE_LIMIT))
        selected = sorted(
            deduplicated.values(),
            key=lambda item: (-item["priority_score"], item["source"]["sheet"], item["source"]["row"]),
        )[:max_evidence]

        for item in selected:
            item.pop("priority_score", None)

        package = {
            "schema_version": "research-evidence/1.0",
            "project": {
                "title": "Demo Health Concept Test Market Research Report",
                "client": "Demo Health Co.",
                "research_type": "new_product_concept_test",
                "language": "zh-CN",
                "business_focus": payload.get("business_focus") or "",
                "audiences": ["Partner", "Consumer", "All"],
                "concepts": [
                    {"concept_id": "C1", "name": "Concept A"},
                    {"concept_id": "C2", "name": "Concept B"},
                ],
            },
            "metric_definitions": [],
            "evidence_registry": selected,
            "reference_report_outline": _reference_outline(reference_ppt),
            "quality": {
                "source_sheet_count": len(workbook.sheetnames),
                "table_count": len(tables),
                "raw_evidence_count": len(all_evidence),
                "deduplicated_evidence_count": len(deduplicated),
                "selected_evidence_count": len(selected),
                "requested_evidence_limit": requested_max_evidence,
                "applied_evidence_limit": max_evidence,
                "truncated": len(deduplicated) > len(selected),
                "warnings": [
                    "Very small bases must not be used as primary evidence.",
                    "Significance letters are retained as source-column references.",
                ],
            },
        }

        def refresh_package():
            used_metric_ids = {item["metric_id"] for item in selected}
            package["metric_definitions"] = [
                metric for metric_id, metric in metric_map.items()
                if metric_id in used_metric_ids
            ]
            package["evidence_registry"] = selected
            package["quality"]["selected_evidence_count"] = len(selected)
            return json.dumps(package, ensure_ascii=False, separators=(",", ":"))

        serialized = refresh_package()
        initial_selected_count = len(selected)
        while len(serialized) > DIFY_SAFE_OUTPUT_CHARS and len(selected) > 100:
            del selected[max(100, len(selected) - 25):]
            serialized = refresh_package()
        package["quality"]["length_trimmed"] = len(selected) < initial_selected_count
        package["quality"]["serialized_character_count"] = len(serialized)
        package["quality"]["truncated"] = len(deduplicated) > len(selected)
        return {
            "status": "success",
            "schema_version": "research-evidence/1.0",
            "research_evidence": package,
        }
    finally:
        if workbook is not None:
            workbook.close()
        local_file.unlink(missing_ok=True)
